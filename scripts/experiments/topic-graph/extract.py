#!/usr/bin/env python3
"""
Stage B — dedup-first extraction + embed-input assembly.

Reads nodes.jsonl (Stage A) and produces extract.jsonl: for every node that will
enter the graph, the text we will actually embed (`embed_input`).

Two principles forced by the real data:
  1. DEDUP FIRST. Group file nodes by Stage-A byte_fp; extract/embed ONE primary
     per group. Non-primaries get `dup_of` and no embed_input (~46% of files).
  2. Extraction differs by type (the matrix). Data files (.json/.yaml dumps) are
     represented LIGHTLY — filename + top-level keys — option (a): the collection
     node carries their topic, individual records stay searchable via their light
     vector, and cosine dedup collapses the near-identical ones in Stage C.

Ports to Swift extraction services (PDFKit/Vision swap in for pdftotext/textutil);
the matrix logic is the portable part.
"""
from __future__ import annotations

import json
import os
import re
import subprocess
from collections import Counter
from pathlib import Path

try:
    import yaml
except Exception:
    yaml = None
try:
    import docx  # python-docx
except Exception:
    docx = None
try:
    from pptx import Presentation
except Exception:
    Presentation = None
try:
    import openpyxl
except Exception:
    openpyxl = None
try:
    from bs4 import BeautifulSoup
except Exception:
    BeautifulSoup = None

HERE = Path(__file__).parent
MAX_CHARS = 8000          # ~2k tokens; bge-m3 handles 8k tokens, plenty for a topic
CODE_HEAD_LINES = 200     # how many leading lines of source to scan for symbols
COLLECTION_SAMPLE = 40    # member titles to sample into a collection's input
COLLECTION_SNIPPETS = 4   # content members to pull a snippet from

SYMBOL_RE = re.compile(
    r"^\s*(?:export\s+)?(?:public\s+|private\s+|internal\s+|static\s+|async\s+)*"
    r"(?:def|class|func|struct|enum|protocol|interface|fn|function|type|trait|impl)\s+"
    r"([A-Za-z_][A-Za-z0-9_]*)", re.MULTILINE)
COMMENT_RE = re.compile(r'^\s*(?://|#|/\*|\*|"""|<!--)')


def clip(s: str, n: int = MAX_CHARS) -> str:
    s = re.sub(r"[ \t]+", " ", (s or "")).strip()
    return s[:n]


def run(cmd: list[str], timeout: int = 20) -> str:
    try:
        r = subprocess.run(cmd, capture_output=True, timeout=timeout, text=True)
        return r.stdout or ""
    except Exception:
        return ""


# ---------------------------------------------------------------- file extractors

def extract_pdf(path: str) -> str:
    return run(["pdftotext", "-l", "12", "-q", path, "-"])


def extract_docx(path: str) -> str:
    if docx:
        try:
            return "\n".join(p.text for p in docx.Document(path).paragraphs)
        except Exception:
            pass
    return run(["textutil", "-convert", "txt", "-stdout", path])


def extract_pptx(path: str) -> str:
    if not Presentation:
        return ""
    try:
        out = []
        for slide in Presentation(path).slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    out.append(shape.text_frame.text)
        return "\n".join(out)
    except Exception:
        return ""


def extract_xlsx(path: str) -> str:
    if not openpyxl:
        return ""
    try:
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        out = []
        for ws in wb.worksheets[:4]:
            out.append(f"[sheet: {ws.title}]")
            for i, row in enumerate(ws.iter_rows(values_only=True)):
                if i >= 6:
                    break
                out.append(" | ".join(str(c) for c in row if c is not None))
        return "\n".join(out)
    except Exception:
        return ""


def extract_html(path: str) -> str:
    try:
        raw = Path(path).read_text(errors="ignore")
    except OSError:
        return ""
    if BeautifulSoup:
        try:
            return BeautifulSoup(raw, "html.parser").get_text(" ")
        except Exception:
            pass
    return re.sub(r"<[^>]+>", " ", raw)


def extract_text(path: str) -> str:
    try:
        return Path(path).read_text(errors="ignore")
    except OSError:
        return ""


def extract_rtf_doc(path: str) -> str:
    return run(["textutil", "-convert", "txt", "-stdout", path])


def extract_json_light(path: str) -> str:
    """Option (a): represent a data file by its SHAPE, not its content — top-level
    keys (or keys of the first record). Cheap, and cosine collapses near-dups."""
    try:
        raw = Path(path).read_text(errors="ignore")[:200_000]
        obj = json.loads(raw)
    except Exception:
        return ""
    if isinstance(obj, dict):
        keys = list(obj.keys())
    elif isinstance(obj, list) and obj and isinstance(obj[0], dict):
        keys = list(obj[0].keys())
    else:
        keys = []
    return "fields: " + ", ".join(map(str, keys[:40]))


def extract_yaml_light(path: str) -> str:
    if not yaml:
        return ""
    try:
        obj = yaml.safe_load(Path(path).read_text(errors="ignore")[:200_000])
        if isinstance(obj, dict):
            return "keys: " + ", ".join(map(str, list(obj.keys())[:40]))
    except Exception:
        pass
    return ""


def extract_source(path: str) -> str:
    """Leading comment block + symbol names — the gist of a source file."""
    txt = extract_text(path)
    if not txt:
        return ""
    lines = txt.splitlines()[:CODE_HEAD_LINES]
    comments = [ln.strip() for ln in lines[:30] if COMMENT_RE.match(ln)]
    symbols = list(dict.fromkeys(SYMBOL_RE.findall("\n".join(lines))))  # dedup, keep order
    sym = ", ".join(symbols)
    return " ".join(filter(None, [" ".join(comments), f"symbols: {sym}" if sym else ""]))


DATA_EXTS = {"json", "yaml", "yml"}
SOURCE_EXTS = {"swift", "js", "ts", "tsx", "jsx", "py", "java", "c", "cpp", "h",
               "hpp", "rb", "go", "rs", "kt", "m", "mm", "php", "lua", "r", "sh", "sql"}


def extract_file(node: dict) -> tuple[str, bool]:
    """Return (text, is_data). is_data flags lightly-represented records."""
    path, ext = node["path"], (node.get("ext") or "")
    if ext == "pdf":
        return extract_pdf(path), False
    if ext in ("docx", "doc"):
        return (extract_docx(path) if ext == "docx" else extract_rtf_doc(path)), False
    if ext in ("pptx", "ppt"):
        return extract_pptx(path), False
    if ext in ("xlsx", "xls", "csv", "tsv", "numbers"):
        return (extract_xlsx(path) if ext in ("xlsx", "xls")
                else extract_text(path)), False
    if ext in ("md", "txt", "rtf"):
        return (extract_rtf_doc(path) if ext == "rtf" else extract_text(path)), False
    if ext in ("html", "htm"):
        return extract_html(path), False
    if ext == "json":
        return extract_json_light(path), True
    if ext in ("yaml", "yml"):
        return extract_yaml_light(path), True
    if ext in SOURCE_EXTS:
        return extract_source(path), False
    return "", False


# ----------------------------------------------------------- project / collection

def find_readme(dirpath: str) -> str:
    try:
        for name in os.listdir(dirpath):
            if name.lower().startswith("readme"):
                return clip(extract_text(os.path.join(dirpath, name)), 4000)
    except OSError:
        pass
    return ""


def manifest_desc(dirpath: str) -> str:
    out = []
    pj = os.path.join(dirpath, "package.json")
    if os.path.exists(pj):
        try:
            d = json.loads(Path(pj).read_text(errors="ignore"))
            out += [d.get("name", ""), d.get("description", ""),
                    " ".join(d.get("keywords", []) or []),
                    " ".join(list((d.get("dependencies") or {}).keys())[:20])]
        except Exception:
            pass
    for tf in ("pyproject.toml", "Cargo.toml"):
        p = os.path.join(dirpath, tf)
        if os.path.exists(p):
            txt = extract_text(p)
            for m in re.findall(r'(?:name|description)\s*=\s*["\']([^"\']+)["\']', txt):
                out.append(m)
    return " ".join(filter(None, out))


def project_input(node: dict) -> str:
    d = node["path"]
    name = Path(d).name
    readme = find_readme(d)
    mani = manifest_desc(d)
    langs = ", ".join(node.get("languages") or [])
    try:
        tree = ", ".join([x for x in sorted(os.listdir(d)) if not x.startswith(".")][:25])
    except OSError:
        tree = ""
    return clip(f"project: {name}. languages: {langs}. {mani}. "
                f"contents: {tree}. {readme}")


def collection_input(node: dict, members: list[dict]) -> str:
    name = Path(node["path"]).name
    titles = [Path(m["path"]).stem for m in members
              if not Path(m["path"]).name.startswith(".")][:COLLECTION_SAMPLE]
    snippets = []
    for m in members:
        if len(snippets) >= COLLECTION_SNIPPETS:
            break
        if m.get("ext") in ("md", "txt"):
            t = extract_text(m["path"])
            if t.strip():
                snippets.append(clip(t, 400))
    return clip(f"collection: {name}. items: {', '.join(titles)}. "
                + " ".join(snippets))


# --------------------------------------------------------------------------- main

def main():
    nodes = [json.loads(l) for l in open(HERE / "nodes.jsonl")]
    by_id = {n["id"]: n for n in nodes}

    # ---- dedup-first: one primary per byte_fp group (prefer shortest path) ----
    groups: dict[str, list[dict]] = {}
    for n in nodes:
        if n["type"] == "file" and n.get("byte_fp"):
            groups.setdefault(n["byte_fp"], []).append(n)
    dup_of: dict[str, str] = {}
    for fp, grp in groups.items():
        primary = min(grp, key=lambda x: (len(x["path"]), x["path"]))
        for n in grp:
            if n["id"] != primary["id"]:
                dup_of[n["id"]] = primary["id"]

    # members per collection (primaries only, content-bearing)
    coll_members: dict[str, list[dict]] = {}
    for n in nodes:
        if n["type"] == "file" and n.get("collection_id"):
            coll_members.setdefault(n["collection_id"], []).append(n)

    out = open(HERE / "extract.jsonl", "w")
    st = Counter()
    total_chars = 0

    for n in nodes:
        rec = {"id": n["id"], "type": n["type"], "path": n["path"],
               "root": n.get("root"), "parent_id": n.get("parent_id"),
               "collection_id": n.get("collection_id")}

        if n["id"] in dup_of:
            rec["dup_of"] = dup_of[n["id"]]
            out.write(json.dumps(rec) + "\n")
            st["dup_skipped"] += 1
            continue

        if n["type"] == "project":
            text = project_input(n)
            rec["namespace"] = "project"
            st["project"] += 1
        elif n["type"] == "collection":
            text = collection_input(n, coll_members.get(n["id"], []))
            rec["namespace"] = "project"   # collections cluster WITH projects
            st["collection"] += 1
        elif n["type"] == "file":
            if not n.get("content_bearing"):
                rec["namespace"] = "meta"  # by-type lane, no topic embedding
                rec["embed_input"] = Path(n["path"]).name
                out.write(json.dumps(rec) + "\n")
                st["meta_only"] += 1
                continue
            body, is_data = extract_file(n)
            text = clip(f"{Path(n['path']).stem}. {body}")
            rec["namespace"] = "data" if is_data else "doc"
            rec["is_data"] = is_data
            st["data" if is_data else "doc"] += 1
        else:
            continue

        text = clip(text)
        rec["embed_input"] = text
        rec["chars"] = len(text)
        total_chars += len(text)
        if len(text) < 12:
            st["empty_extract"] += 1
        out.write(json.dumps(rec) + "\n")

    out.close()
    report(st, total_chars, len(nodes))


def report(st: Counter, total_chars: int, n_nodes: int):
    line = "=" * 60
    to_embed = st["project"] + st["collection"] + st["doc"] + st["data"]
    print(line)
    print("STAGE B — EXTRACTION (dedup-first, option (a) for data)")
    print(line)
    print(f"input nodes:            {n_nodes}")
    print(f"dup-skipped (no embed): {st['dup_skipped']}")
    print(f"meta-only (by-type):    {st['meta_only']}")
    print(line)
    print("to embed:")
    print(f"  projects:             {st['project']}")
    print(f"  collections:          {st['collection']}")
    print(f"  documents:            {st['doc']}")
    print(f"  data (light/json):    {st['data']}")
    print(f"  TOTAL embed nodes:    {to_embed}")
    print(f"  empty extractions:    {st['empty_extract']}")
    print(line)
    est_tokens = total_chars // 4
    print(f"total embed chars:      {total_chars:,}  (~{est_tokens:,} tokens)")
    print(f"extract.jsonl written. (was {n_nodes} nodes → {to_embed} to embed)")


if __name__ == "__main__":
    main()
